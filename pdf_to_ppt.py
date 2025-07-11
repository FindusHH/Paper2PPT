"""Utility helpers used by the Streamlit front end.

This module handles extracting content from PDF files,
interacting with Azure OpenAI to create summaries and titles,
and building the final PowerPoint presentation."""
import io
from pathlib import Path
from typing import List
import base64
import json

from langdetect import detect

# Location of the text files containing the prompts

# Path to the system prompt used for summarization
PROMPT_PATH = Path(__file__).resolve().parent / "prompts" / "summarize.txt"
# Path to the image relevance evaluation prompt
IMAGE_PROMPT_PATH = Path(__file__).resolve().parent / "prompts" / "image_eval.txt"
# Path to the slide title prompt
TITLE_PROMPT_PATH = Path(__file__).resolve().parent / "prompts" / "title.txt"

# Settings file controlling language options and formatting
SETTINGS_FILE = Path(__file__).resolve().parent / "settings.json"


# Default values used when settings.json is missing

DEFAULT_SETTINGS = {
    "languages": {
        "English": "en",
        "German": "de",
        "Spanish": "es",
        "Chinese": "zh",
    },
    "font_size": 24,
    "max_words_per_bullet": 10,
    "max_words_title": 4,

    "min_image_score": 5,
    "pages_per_slide": 1,

}


def load_settings() -> dict:

    """Read settings from settings.json or use defaults."""

    if SETTINGS_FILE.exists():
        with open(SETTINGS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    return DEFAULT_SETTINGS



def load_prompt(path: Path = PROMPT_PATH) -> str:
    """Load the system prompt from the prompts directory."""
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read().strip()
    except FileNotFoundError:
        # Fallback prompt if the file does not exist
        return (
            "Summarize the following text into at most 5 concise bullet points."
            " Respond in {language}."
        )


def save_prompt(content: str, path: Path = PROMPT_PATH) -> None:
    """Persist the system prompt to disk."""
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)


SYSTEM_PROMPT = load_prompt()
IMAGE_PROMPT = load_prompt(IMAGE_PROMPT_PATH)
TITLE_PROMPT = load_prompt(TITLE_PROMPT_PATH)

SETTINGS = load_settings()


import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_pages(pdf_path: str):
    """Extract text and images from each page of a PDF."""
    # Open the PDF with PyMuPDF
    doc = fitz.open(pdf_path)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        images = []
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            ext = base_image["ext"]
            images.append((image_bytes, ext))
        yield page_num + 1, text, images

    # Close the document to free resources
    doc.close()


def detect_pdf_language(pdf_path: str) -> str:
    """Detect predominant language of the PDF text."""
    text_snippets = []
    for _, text, _ in extract_pages(pdf_path):
        text_snippets.append(text)
        if len(" ".join(text_snippets)) > 1000:
            break
    try:
        return detect(" ".join(text_snippets))
    except Exception:
        return "en"


def create_slide(prs: Presentation, title: str, bullets: List[str], images: List[bytes]):
    """Add a slide with a title, bullet points and images."""
    # Use the "Title and Content" layout
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title


    body_placeholder = slide.shapes.placeholders[1]
    body_placeholder.left = Inches(0.5)
    body_placeholder.top = Inches(1.0)
    body_placeholder.height = Inches(4.0)
    if images:
        body_placeholder.width = Inches(5)
        pic_left = Inches(5.6)
    else:
        body_placeholder.width = Inches(9)
        pic_left = None
    body = body_placeholder.text_frame
    # Remove any existing text from the placeholder
    body.clear()

    for point in bullets:
        p = body.add_paragraph()
        p.text = point
        p.level = 0

        size = min(SETTINGS.get("font_size", 24), 32)
        p.font.size = Pt(size)
    if images:
        img_bytes, ext = images[0]
        image_stream = io.BytesIO(img_bytes)
        slide.shapes.add_picture(image_stream, pic_left, Inches(1.5), height=Inches(4))


def _add_bullet_slides(prs: Presentation, title: str, bullets: List[str], images: List[bytes]):
    """Create one or more slides ensuring bullet lists fit."""

    # Only five bullets fit on a single slide
    MAX_BULLETS = 5
    # Guard against unexpected None values from upstream code
    if bullets is None:
        bullets = []


    for idx in range(0, len(bullets), MAX_BULLETS):
        group = bullets[idx : idx + MAX_BULLETS]
        slide_title = title if idx == 0 else f"{title} (cont.)"
        create_slide(prs, slide_title, group, images if idx == 0 else [])


def save_presentation(sections, output_path: str):
    """Write all slides to a PowerPoint file."""

    prs = Presentation()
    # Add each section of content as one or more slides

    for title, bullets, images in sections:
        _add_bullet_slides(prs, title, bullets, images)
    # Finally write the presentation to disk

    prs.save(output_path)


from openai import AzureOpenAI


def summarize_text(
    text: str,
    client: AzureOpenAI,
    deployment: str,
    *,
    language: str = "",
    max_tokens: int = 256,
) -> List[str]:

    """Use Azure OpenAI to summarize text into bullet points."""
    system_prompt = SYSTEM_PROMPT

    # Insert the configured word limit into the prompt if needed

    max_words = SETTINGS.get("max_words_per_bullet", 10)
    if "{max_words}" in system_prompt:
        system_prompt = system_prompt.replace("{max_words}", str(max_words))
    if "{language}" in system_prompt:
        system_prompt = system_prompt.replace("{language}", language or "the original language")
    elif language:
        system_prompt = f"{system_prompt}\nRespond in {language}."

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": text},
    ]

    try:
        response = client.chat.completions.create(
            model=deployment,
            messages=messages,
            max_tokens=max_tokens,
        )
        content = response.choices[0].message.content or ""

        bullets = [line.lstrip("- ").strip() for line in content.splitlines() if line]
        trimmed = []
        # Truncate each bullet to the configured word limit
        for b in bullets:
            words = b.split()
            if len(words) > max_words:
                words = words[:max_words]
            trimmed.append(" ".join(words))
        return trimmed
    except Exception:
        # In case of API failure, return an empty list instead of None
        return []


def generate_title(
    text: str,
    client: AzureOpenAI,
    deployment: str,
    *,
    language: str = "",
    max_tokens: int = 16,
) -> str:
    """Generate a short slide title."""


    max_words = SETTINGS.get("max_words_title", 4)
    prompt = TITLE_PROMPT
    # Replace placeholders in the title prompt
    if "{max_words}" in prompt:
        prompt = prompt.replace("{max_words}", str(max_words))
    if "{language}" in prompt:
        prompt = prompt.replace("{language}", language or "the original language")
    elif language:
        prompt = f"{prompt}\nRespond in {language}."
    # Prepare the conversation for the chat completion call

    messages = [
        {"role": "system", "content": prompt},
        {"role": "user", "content": text},
    ]

    try:
        response = client.chat.completions.create(
            model=deployment,
            messages=messages,
            max_tokens=max_tokens,
        )
        content = response.choices[0].message.content
        if content:
            text = content.strip().strip('"')
            return text
    except Exception:
        # Fall back to empty title on API error
        pass
    return ""



def evaluate_image_relevance(
    page_text: str,
    image: bytes,
    ext: str,
    client: AzureOpenAI,
    deployment: str,
    *,
    max_tokens: int = 8,
) -> float:
    """Return an image relevance score between 0 and 10."""

    # We ask the language model whether the image clarifies the given
    # page text. The API expects a data URL for the image content.

    b64 = base64.b64encode(image).decode("utf-8")
    mime = f"data:image/{ext};base64,{b64}"
    # Combine the page text and the image into a single chat request
    messages = [
        {"role": "system", "content": IMAGE_PROMPT},
        {
            "role": "user",
            "content": [
                {"type": "text", "text": page_text},
                {"type": "image_url", "image_url": {"url": mime}},
            ],
        },
    ]
    try:
        response = client.chat.completions.create(
            model=deployment,
            messages=messages,
            max_tokens=max_tokens,
        )
        answer = response.choices[0].message.content.strip()
        return float(answer)
    except Exception:
        return 0.0



def pdf_to_ppt(
    pdf_path: str,
    output_path: str,
    client: AzureOpenAI,
    deployment: str,
    *,
    language: str = "",

    pages_per_slide: int = 1,
    progress_callback=None,
) -> None:
    """Convert a PDF document to a summarized PowerPoint file.

    ``pages_per_slide`` controls how many PDF pages are combined before
    generating a single slide. The highest scoring image from that group
    is used if its relevance surpasses the configured minimum score.
    """


    # Reload settings in case they were changed
    global SETTINGS
    SETTINGS = load_settings()

    # Collect (title, bullets, [image]) tuples for each group of pages

    sections = []
    # Minimum relevance score an image must achieve to be used
    min_score = SETTINGS.get("min_image_score", 5)


    page_data = list(extract_pages(pdf_path))
    total_groups = (len(page_data) + pages_per_slide - 1) // pages_per_slide

    for group_idx in range(total_groups):
        group = page_data[group_idx * pages_per_slide : (group_idx + 1) * pages_per_slide]
        combined_text = "\n".join(p[1] for p in group)
        group_images = [img for p in group for img in p[2]]
        if progress_callback:
            progress_callback(group_idx + 1, total_groups, f"Part {group_idx + 1}/{total_groups}")
        title = generate_title(combined_text, client, deployment, language=language)
        bullets = summarize_text(combined_text, client, deployment, language=language)

        best_img = None
        best_score = -1.0
        # Evaluate all images and keep the highest scoring one
        for img_bytes, ext in group_images:
            score = evaluate_image_relevance(combined_text, img_bytes, ext, client, deployment)
            if score > best_score:
                best_score = score
                best_img = (img_bytes, ext)

        relevant_images = [best_img] if best_img and best_score >= min_score else []


        sections.append((title, bullets, relevant_images))
    # Write all collected slides to the output file
    save_presentation(sections, output_path)
    if progress_callback:
        progress_callback(total_groups, total_groups, "Completed")


